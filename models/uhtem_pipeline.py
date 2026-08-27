# uhtem_pipeline.py - Unified-Hybrid Text Extraction Model (UHTEM) Engine
import os
import sys
import time
import psutil
from typing import Dict, Any, List, Tuple
from PIL import Image

# Enforce stable settings for Paddle/Windows environments to prevent crash
os.environ["FLAGS_use_onednn"] = "0"
os.environ["FLAGS_use_dx12"] = "0"
os.environ["FLAGS_use_mkldnn"] = "0"
os.environ["PADDLE_PDX_ENABLE_MKLDNN_BYDEFAULT"] = "0"

class UHTEMEngine:
    """
    Unified-Hybrid Text Extraction Model (UHTEM) Engine.
    Combines the lightweight OCR speed of PaddleOCR with the layout-aware coordinate system of LayoutLMv3.
    Incorporates an Adaptive Router to run entirely on CPU/lightweight devices without GPU over-reliance.

    Supports: PDF (digital & scanned), DOCX, PPTX, and image files.
    """
    def __init__(self, use_gpu: str = "auto", low_resource_mode: bool = True):
        self.low_resource_mode = low_resource_mode
        self.paddle_ocr = None
        self.layoutlm_processor = None
        self.layoutlm_model = None

        # Automatically detect NVIDIA CUDA capability at runtime via nvidia-smi
        if use_gpu == "auto":
            try:
                import subprocess
                subprocess.run(["nvidia-smi"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
                self.use_gpu = True
            except Exception:
                self.use_gpu = False
        else:
            self.use_gpu = bool(use_gpu)

        print(f"[UHTEM] Initialized Engine (use_gpu={self.use_gpu}, low_resource_mode={self.low_resource_mode})")

    def warmup(self):
        """Pre-loads PaddleOCR and runs a dummy inference to compile C++ backends."""
        try:
            self._init_paddle()
            # Create a tiny temp image to compile C++ backends
            from PIL import Image
            temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uhtem_temp_cache")
            os.makedirs(temp_dir, exist_ok=True)
            temp_img_path = os.path.join(temp_dir, "warmup_dummy.png")
            
            img = Image.new("RGB", (100, 100), color="white")
            img.save(temp_img_path)
            
            try:
                self.paddle_ocr.predict(temp_img_path)
            except Exception:
                pass
                
            try:
                os.remove(temp_img_path)
            except OSError:
                pass
            print("[UHTEM] Background warmup completed successfully.")
        except Exception as e:
            print(f"[WARNING] UHTEM warmup failed: {e}")

    def _init_paddle(self):
        """Lazy load PaddleOCR only when scanned documents require OCR, saving initial memory."""
        if self.paddle_ocr is None:
            print("[UHTEM] Lazy-initializing CPU-optimized PaddleOCR mobile backbone...")
            try:
                from paddleocr import PaddleOCR
                # use_angle_cls=True detects oriented text. lang='en' handles standard text.
                # enforce CPU execution if use_gpu=False to prevent VRAM allocation spikes.
                try:
                    self.paddle_ocr = PaddleOCR(use_angle_cls=True, lang='en', use_gpu=self.use_gpu, text_det_limit_side_len=320)
                except ValueError:
                    print("[UHTEM] 'use_gpu' argument not supported in this version. Falling back to default initialization.")
                    self.paddle_ocr = PaddleOCR(use_angle_cls=True, lang='en', text_det_limit_side_len=320)
                print("[UHTEM] PaddleOCR backbone initialized successfully.")
            except ImportError:
                print("[ERROR] PaddleOCR is not installed. Fallback to native text extraction only.")
                raise

    def _init_layoutlmv3(self):
        """Lazy load LayoutLMv3 components if full transformer classification is requested (non-low_resource_mode)."""
        if not self.low_resource_mode and self.layoutlm_processor is None:
            print("[UHTEM] Loading heavy LayoutLMv3 sequence model and processor...")
            try:
                import torch
                from transformers import LayoutLMv3Processor, LayoutLMv3ForSequenceClassification
                self.layoutlm_processor = LayoutLMv3Processor.from_pretrained("microsoft/layoutlmv3-base", apply_ocr=False)
                self.layoutlm_model = LayoutLMv3ForSequenceClassification.from_pretrained("microsoft/layoutlmv3-base")
                if self.use_gpu and torch.cuda.is_available():
                    self.layoutlm_model = self.layoutlm_model.to("cuda")
                print("[UHTEM] LayoutLMv3 model loaded successfully.")
            except Exception as e:
                print(f"[WARNING] LayoutLMv3 loading failed: {e}. Falling back to layout dictionary extraction.")
                self.low_resource_mode = True

    def normalize_box(self, box: List[float], width: float, height: float) -> List[int]:
        """Normalizes bounding boxes to LayoutLMv3's standard [0, 1000] scale."""
        x0, y0, x1, y1 = box
        # Scale to 0-1000 coordinate grid
        nx0 = max(0, min(1000, int(1000 * x0 / width)))
        ny0 = max(0, min(1000, int(1000 * y0 / height)))
        nx1 = max(0, min(1000, int(1000 * x1 / width)))
        ny1 = max(0, min(1000, int(1000 * y1 / height)))
        return [nx0, ny0, nx1, ny1]

    def _extract_digital_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extracts words, normalized coordinates, and rendered images directly from digital PDF.
        Highly optimized: 0% GPU utilization, runs purely on CPU in milliseconds.
        """
        import fitz  # PyMuPDF
        print(f"[UHTEM-Router] Routing to DIRECT DIGITAL PARSER for: {os.path.basename(pdf_path)}")
        
        doc = fitz.open(pdf_path)
        pages_data = []

        for page_idx in range(len(doc)):
            page = doc.load_page(page_idx)
            width = page.rect.width
            height = page.rect.height

            # Render high-resolution page image for visual features/LayoutLMv3 input
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            words_list = []
            boxes_list = []

            # page.get_text("words") returns list of tuples: (x0, y0, x1, y1, "word", block_no, line_no, word_no)
            raw_words = page.get_text("words")
            
            for w in raw_words:
                word_text = w[4]
                # Filter out white spaces
                if not word_text.strip():
                    continue
                
                # Physical coordinates
                box = [w[0], w[1], w[2], w[3]]
                # Normalize box to LayoutLMv3 space
                norm_box = self.normalize_box(box, width, height)

                words_list.append(word_text)
                boxes_list.append(norm_box)

            pages_data.append({
                "page_number": page_idx + 1,
                "extraction_method": "PyMuPDF-DigitalDirect",
                "width": width,
                "height": height,
                "image": img,
                "words": words_list,
                "boxes": boxes_list,
                "device_used": "CPU"
            })
            
        doc.close()
        return pages_data

    def _extract_scanned_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from scanned files or pure images using lightweight mobile PaddleOCR.
        Bypasses native extraction, performing high-fidelity OCR on CPU (or GPU if specified).
        """
        self._init_paddle()
        import fitz
        print(f"[UHTEM-Router] Routing to LIGHTWEIGHT PADDLEOCR ENGINE for: {os.path.basename(file_path)}")

        image_paths = []
        is_pdf = file_path.lower().endswith(".pdf")
        temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uhtem_temp_cache")
        os.makedirs(temp_dir, exist_ok=True)

        # Step 1: Render PDF pages to images if input is a scanned PDF
        if is_pdf:
            doc = fitz.open(file_path)
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=150)
                img_path = os.path.join(temp_dir, f"temp_page_{i}.png")
                pix.save(img_path)
                image_paths.append((i + 1, img_path))
            doc.close()
        else:
            # Single image input
            image_paths.append((1, file_path))

        pages_data = []

        # Step 2: Run lightweight OCR on each page image
        for page_num, img_path in image_paths:
            img = Image.open(img_path)
            width, height = img.size

            # Run PaddleOCR predict directly to bypass broken internal ocr() method
            try:
                result = self.paddle_ocr.predict(img_path)
            except Exception as e:
                print(f"[WARNING] OCR execution error: {e}")
                result = None

            words_list = []
            boxes_list = []

            # Parse results
            if result and len(result) > 0:
                res = result[0]
                
                # Format A: Dictionary-based (PaddleOCR 3.5.0 / PaddleX)
                if hasattr(res, 'get') and 'rec_texts' in res:
                    texts = res.get('rec_texts', [])
                    # Look for polygon boundaries or boxes
                    raw_boxes = res.get('dt_polys', res.get('dt_boxes', []))
                    
                    for idx, text in enumerate(texts):
                        x0, y0, x1, y1 = 0, 0, width, height
                        if idx < len(raw_boxes):
                            coords = raw_boxes[idx]
                            # Handle different coordinate structures safely
                            if isinstance(coords, list) and len(coords) > 0:
                                if isinstance(coords[0], list):
                                    x0 = min(pt[0] for pt in coords)
                                    y0 = min(pt[1] for pt in coords)
                                    x1 = max(pt[0] for pt in coords)
                                    y1 = max(pt[1] for pt in coords)
                                else:
                                    x0, y0, x1, y1 = coords
                        
                        norm_box = self.normalize_box([x0, y0, x1, y1], width, height)
                        sub_words = text.split()
                        for sw in sub_words:
                            words_list.append(sw)
                            boxes_list.append(norm_box)
                            
                # Format B: Legacy List-based ([[[coords], (text, conf)]])
                elif isinstance(res, list):
                    for line in res:
                        if isinstance(line, list) and len(line) >= 2:
                            coords = line[0]
                            text, conf = line[1]
                            
                            x0 = min(pt[0] for pt in coords)
                            y0 = min(pt[1] for pt in coords)
                            x1 = max(pt[0] for pt in coords)
                            y1 = max(pt[1] for pt in coords)

                            norm_box = self.normalize_box([x0, y0, x1, y1], width, height)
                            sub_words = text.split()
                            for sw in sub_words:
                                words_list.append(sw)
                                boxes_list.append(norm_box)

            pages_data.append({
                "page_number": page_num,
                "extraction_method": "PaddleOCR-LightweightMobile",
                "width": width,
                "height": height,
                "image": img,
                "words": words_list,
                "boxes": boxes_list,
                "device_used": "GPU" if self.use_gpu else "CPU"
            })

        # Cleanup rendered temp images
        if is_pdf:
            for _, path in image_paths:
                try:
                    os.remove(path)
                except OSError:
                    pass
            try:
                os.rmdir(temp_dir)
            except OSError:
                pass

        return pages_data

    def _extract_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from DOCX (Word) files using python-docx.
        Extracts paragraphs and tables natively, and extracts inline/embedded
        images to run them through PaddleOCR for full hybrid document parsing.
        """
        import re
        import zipfile
        from docx import Document
        print(f"[UHTEM-Router] Routing to DOCX PARSER (with Hybrid OCR support) for: {os.path.basename(file_path)}")

        doc = Document(file_path)
        words_list = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                words_list.extend(text.split())

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        words_list.extend(text.split())

        # ── Embedded Image Extraction & OCR ─────────────────────
        ocr_ran = False
        temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uhtem_temp_cache")
        
        try:
            # Check for embedded media in the zip file
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(file_path, 'r') as zf:
                    # List all media files
                    media_files = [f for f in zf.namelist() if f.startswith('word/media/')]
                    # Sort files sequentially to process in chronological layout order
                    media_files.sort(key=lambda x: [int(s) if s.isdigit() else s for s in re.split(r'(\d+)', x)])

                    if media_files:
                        os.makedirs(temp_dir, exist_ok=True)
                        self._init_paddle()  # Lazy load PaddleOCR
                        
                        for idx, media_file in enumerate(media_files):
                            print(f"[UHTEM-DOCX] Running OCR on embedded image: {media_file}")
                            img_data = zf.read(media_file)
                            ext = os.path.splitext(media_file)[1] or ".png"
                            temp_img_path = os.path.join(temp_dir, f"docx_temp_img_{idx}{ext}")
                            
                            # Downscale image to a maximum width of 512px to speed up CPU OCR processing
                            try:
                                import io
                                from PIL import Image
                                img = Image.open(io.BytesIO(img_data))
                                # Convert mode if saving as JPEG
                                if ext.lower() in ['.jpg', '.jpeg'] and img.mode in ("RGBA", "P"):
                                    img = img.convert("RGB")
                                width, height = img.size
                                if width > 512:
                                    ratio = 512.0 / width
                                    new_height = int(height * ratio)
                                    try:
                                        img = img.resize((512, new_height), Image.Resampling.LANCZOS)
                                    except AttributeError:
                                        img = img.resize((512, new_height), Image.ANTIALIAS)
                                img.save(temp_img_path)
                            except Exception as img_err:
                                print(f"[WARNING] Image preprocessing failed for {media_file}: {img_err}. Writing raw bytes.")
                                with open(temp_img_path, 'wb') as f:
                                    f.write(img_data)
                            
                            # Run PaddleOCR predict directly to bypass broken internal ocr() method
                            try:
                                result = self.paddle_ocr.predict(temp_img_path)
                            except Exception as e:
                                print(f"[WARNING] DOCX OCR failed on image {media_file}: {e}")
                                result = None
 
                            if result and len(result) > 0:
                                ocr_ran = True
                                res = result[0]
                                
                                # Format A: Dictionary-based
                                if hasattr(res, 'get') and 'rec_texts' in res:
                                    for text in res.get('rec_texts', []):
                                        if text.strip():
                                            words_list.extend(text.split())
                                # Format B: List-based
                                elif isinstance(res, list):
                                    for line in res:
                                        if isinstance(line, list) and len(line) >= 2:
                                            text, conf = line[1]
                                            if text.strip():
                                                words_list.extend(text.split())
                            
                            # Cleanup this temp image immediately
                            try:
                                os.remove(temp_img_path)
                            except OSError:
                                pass
        except Exception as e:
            print(f"[WARNING] DOCX Zip parsing / OCR failed: {e}")
        finally:
            # Ensure cleanup of temp dir
            try:
                if os.path.exists(temp_dir):
                    # Check if empty, then delete
                    if not os.listdir(temp_dir):
                        os.rmdir(temp_dir)
            except OSError:
                pass

        extraction_method = "python-docx-HybridOCR" if ocr_ran else "python-docx-NativeText"
        device_used = ("GPU" if self.use_gpu else "CPU") if ocr_ran else "CPU"

        return [{
            "page_number": 1,
            "extraction_method": extraction_method,
            "width": 0,
            "height": 0,
            "image": None,
            "words": words_list,
            "boxes": [],
            "device_used": device_used
        }]

    def _extract_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from PPTX (PowerPoint) files using python-pptx.
        Iterates through each slide, extracts text from all shapes/tables,
        and runs OCR on slide pictures using PaddleOCR.
        """
        import os
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        print(f"[UHTEM-Router] Routing to PPTX PARSER (with Hybrid OCR support) for: {os.path.basename(file_path)}")

        prs = Presentation(file_path)
        pages_data = []
        temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uhtem_temp_cache")

        for slide_idx, slide in enumerate(prs.slides):
            words_list = []
            ocr_ran = False

            # Extract native text from shapes and tables
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            words_list.extend(text.split())

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                words_list.extend(text.split())

            # OCR on Slide Images
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        # Use default extensions safely
                        ext = ".png"
                        if hasattr(image, 'ext') and image.ext:
                            ext = f".{image.ext}" if not image.ext.startswith('.') else image.ext
                        
                        os.makedirs(temp_dir, exist_ok=True)
                        temp_img_path = os.path.join(
                            temp_dir, f"pptx_temp_slide_{slide_idx}_pic_{len(words_list)}{ext}"
                        )
                        
                        # Downscale image to a maximum width of 512px to speed up CPU OCR processing
                        try:
                            import io
                            from PIL import Image
                            img = Image.open(io.BytesIO(img_bytes))
                            if ext.lower() in ['.jpg', '.jpeg'] and img.mode in ("RGBA", "P"):
                                img = img.convert("RGB")
                            width, height = img.size
                            if width > 512:
                                ratio = 512.0 / width
                                new_height = int(height * ratio)
                                try:
                                    img = img.resize((512, new_height), Image.Resampling.LANCZOS)
                                except AttributeError:
                                    img = img.resize((512, new_height), Image.ANTIALIAS)
                            img.save(temp_img_path)
                        except Exception as img_err:
                            print(f"[WARNING] PPTX Image preprocessing failed: {img_err}. Writing raw bytes.")
                            with open(temp_img_path, 'wb') as f:
                                f.write(img_bytes)
 
                        self._init_paddle()  # Lazy load PaddleOCR
                        print(f"[UHTEM-PPTX] Running OCR on embedded picture inside slide {slide_idx + 1}")
                        
                        # Run PaddleOCR predict directly to bypass broken internal ocr() method
                        try:
                            result = self.paddle_ocr.predict(temp_img_path)
                        except Exception as e:
                            print(f"[WARNING] PPTX OCR failed on shape: {e}")
                            result = None
 
                        if result and len(result) > 0:
                            ocr_ran = True
                            res = result[0]
                            
                            # Format A: Dictionary-based
                            if hasattr(res, 'get') and 'rec_texts' in res:
                                for text in res.get('rec_texts', []):
                                    if text.strip():
                                        words_list.extend(text.split())
                            # Format B: List-based
                            elif isinstance(res, list):
                                for line in res:
                                    if isinstance(line, list) and len(line) >= 2:
                                        text, conf = line[1]
                                        if text.strip():
                                            words_list.extend(text.split())
 
                        # Cleanup temp image
                        try:
                            os.remove(temp_img_path)
                        except OSError:
                            pass
                    except Exception as e:
                        print(f"[WARNING] PPTX shape image extraction failed: {e}")

            extraction_method = "python-pptx-HybridOCR" if ocr_ran else "python-pptx-NativeText"
            device_used = ("GPU" if self.use_gpu else "CPU") if ocr_ran else "CPU"

            pages_data.append({
                "page_number": slide_idx + 1,
                "extraction_method": extraction_method,
                "width": 0,
                "height": 0,
                "image": None,
                "words": words_list,
                "boxes": [],
                "device_used": device_used
            })

        # Ensure cleanup of temp dir if empty
        try:
            if os.path.exists(temp_dir) and not os.listdir(temp_dir):
                os.rmdir(temp_dir)
        except OSError:
            pass

        return pages_data

    def extract(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Main entry point for hybrid text extraction.
        Evaluates input type and dynamically routes it to optimize memory and speed.
        Supports: PDF, DOCX, PPTX, and image files.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        start_time = time.time()
        start_memory = psutil.Process(os.getpid()).memory_info().rss / (1024 * 1024) # in MB

        ext = os.path.splitext(file_path)[1].lower()
        pages_data = []

        # ── Adaptive Routing Decision ──────────────────────────
        if ext == ".docx":
            pages_data = self._extract_docx(file_path)

        elif ext == ".pptx":
            pages_data = self._extract_pptx(file_path)

        elif ext == ".pdf":
            # Check if PDF contains selectable digital text
            import fitz
            doc = fitz.open(file_path)
            has_digital_text = False
            for page in doc:
                if page.get_text().strip():
                    has_digital_text = True
                    break
            doc.close()

            if has_digital_text:
                pages_data = self._extract_digital_pdf(file_path)
            else:
                pages_data = self._extract_scanned_file(file_path)
        else:
            # Assume scanned image input
            pages_data = self._extract_scanned_file(file_path)

        # ── LayoutLMv3 Sequence Processing (Only if low_resource_mode is disabled) ──
        if not self.low_resource_mode and len(pages_data) > 0:
            self._init_layoutlmv3()
            for page in pages_data:
                if page.get("image") is None:
                    continue  # Skip DOCX/PPTX pages (no image available)
                print(f"[UHTEM-LayoutLMv3] Encoding layout tokens for page {page['page_number']}...")
                try:
                    # Tokenize and encode words + coordinates
                    encoding = self.layoutlm_processor(
                        page["image"],
                        text=page["words"],
                        boxes=page["boxes"],
                        return_tensors="pt"
                    )
                    # Include standard processor outputs
                    page["layoutlmv3_inputs"] = encoding
                except Exception as e:
                    print(f"[WARNING] LayoutLMv3 tokenization failed for page {page['page_number']}: {e}")

        # Compute efficiency metrics
        end_time = time.time()
        end_memory = psutil.Process(os.getpid()).memory_info().rss / (1024 * 1024)

        for page in pages_data:
            page["metrics"] = {
                "latency_sec": end_time - start_time,
                "memory_used_mb": max(0.0, end_memory - start_memory),
                "peak_system_ram_mb": psutil.virtual_memory().used / (1024 * 1024)
            }

        return pages_data
